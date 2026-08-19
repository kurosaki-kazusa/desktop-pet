# -*- coding: utf-8 -*-
"""基于烽火通信PPT模板生成「AI 桌面宠物」路演 PPT。

用法：python scripts/build_roadshow_ppt.py
输出：docs/AI桌面宠物路演PPT.pptx（不修改模板原文件）
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

sys.stdout.reconfigure(encoding='utf-8')

SRC = r'd:\Documents\VibeCoding课件\projects\docs\烽火通信PPT模板-1.pptx'
OUT = r'd:\Documents\VibeCoding课件\projects\docs\AI桌面宠物路演PPT.pptx'

prs = Presentation(SRC)
slides = list(prs.slides)


def set_run_text(tf, text):
    """替换文本（保留第一段第一个 run 的字体格式），多行用 \\x0b 软换行。"""
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.text = text
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def find_title(slide):
    """返回标题占位符（name 以「标题」开头的 PLACEHOLDER）。"""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name.startswith('标题'):
            return sh
    return None


def find_body_placeholders(slide, top_min=1.0):
    """返回正文类占位符，按 (top, left) 排序，过滤页眉区域。"""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.name.startswith('文本占位符'):
            continue
        top = sh.top
        if top is None or (top / 914400) < top_min:
            continue
        out.append(sh)
    out.sort(key=lambda s: (s.top or 0, s.left or 0))
    return out


def fill_body(slide, subtitle, body):
    """白底/背景色 N 项版式：占位符按高度区分小标题（<1 英寸）与正文（≥1 英寸），
    组内按 (top 量化, left) 排序对齐；subtitle 与 body 一一对应。"""
    small, big = [], []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.name.startswith('文本占位符'):
            continue
        top = sh.top
        if top is None or (top / 914400) < 1.0:
            continue  # 过滤页眉区域小字
        h = sh.height or 0
        (small if h < Inches(1.0) else big).append(sh)

    def key(s):
        top_in = (s.top or 0) / 914400
        return (round(top_in * 20) / 20, (s.left or 0) / 914400)

    small.sort(key=key)
    big.sort(key=key)
    for i, sub in enumerate(subtitle):
        if i < len(small):
            small[i].text_frame.text = sub
    for i, b in enumerate(body):
        if i < len(big):
            big[i].text_frame.text = b


# ============ 1. 封面：加主标题/副标题，改日期 ============
# 封面背景为深蓝（accent1 4472C4 亮度 50%），标题直接使用白色，不继承原 run 格式
s = slides[0]
tb = s.shapes.add_textbox(Inches(1.2), Inches(0.75), Inches(11.0), Inches(1.1))
tf = tb.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'AI 智能桌面宠物'
r.font.size = Pt(44)
r.font.bold = True
r.font.name = '微软雅黑'
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

tb2 = s.shapes.add_textbox(Inches(1.2), Inches(1.95), Inches(11.0), Inches(0.6))
tf2 = tb2.text_frame
tf2.word_wrap = True
r2 = tf2.paragraphs[0].add_run()
r2.text = '基于 Electron 与大模型的陪伴式桌面应用 · 路演汇报'
r2.font.size = Pt(18)
r2.font.name = '微软雅黑'
r2.font.color.rgb = RGBColor(0xD9, 0xE2, 0xF3)

for sh in s.shapes:
    if sh.has_text_frame and sh.name == '矩形 18':
        set_run_text(sh.text_frame, '2026.08')

# ============ 2. 扉页：汇报人信息 ============
s = slides[1]
for sh in s.shapes:
    if sh.has_text_frame and sh.name == '矩形 7':
        set_run_text(sh.text_frame, '烽火通信科技股份有限公司\u3000教学演示组\x0bAI 桌面宠物项目路演\x0b2026年8月')

# ============ 3. 目录（Slide 5，4 项）：主框=章号，副框=章名 ============
s = slides[4]
toc = [
    ('第一章', 'SOP 业务价值'),
    ('第二章', '业务逻辑设计'),
    ('第三章', '技术实现过程'),
    ('第四章', '技术亮点与后续规划'),
]
groups = [sh for sh in s.shapes if sh.shape_type == 6]
# 目录 GROUP 的顺序与视觉位置无关，按组内主标题框 top 排序
def group_top(g):
    for sub in g.shapes:
        if sub.has_text_frame and sub.text_frame.text.strip():
            return sub.top or 0
    return 0
groups.sort(key=group_top)
print(f'目录项数量: {len(groups)}')
for g, (main, sub) in zip(groups, toc):
    boxes = [x for x in g.shapes if x.has_text_frame and x.text_frame.text.strip()]
    boxes.sort(key=lambda b: b.top or 0)
    if boxes:
        set_run_text(boxes[0].text_frame, main)
    if len(boxes) > 1:
        set_run_text(boxes[1].text_frame, sub)

# ============ 4. 章节过渡页（Slide 8-11） ============
chapter_titles = [
    (8, 'SOP 业务价值'),
    (9, '业务逻辑设计'),
    (10, '技术实现过程'),
    (11, '技术亮点与后续规划'),
]
for idx, title in chapter_titles:
    s = slides[idx - 1]
    for sh in s.shapes:
        if sh.has_text_frame and sh.name == 'TextBox 16':
            set_run_text(sh.text_frame, title)
            break

# ============ 5. SOP 业务价值（Slide 15 白底3项） ============
s = slides[14]
find_title(s).text_frame.text = 'SOP 业务价值'
fill_body(s,
    ['教学演示价值', '效率提升价值', '情感陪伴价值'],
    [
        '作为 vibe coding 课程标杆案例，完整演示从需求、原型、编码到打包、AI 接入的标准化全流程，沉淀可复用的 SOP 文档',
        '定时提醒自动调度，常用命令一键复制执行，把重复性桌面操作压缩为一次点击',
        '大模型流式对话让桌宠「活起来」，工作间隙提供陪伴与情绪价值，缓解长时间伏案压力',
    ])

# ============ 6. 项目概况与 SOP 化交付（Slide 13 白底1项） ============
s = slides[12]
find_title(s).text_frame.text = '项目概况：一个桌宠的完整 SOP'
fill_body(s, ['AI 桌面宠物（AI Desk Pet）—— Electron 悬浮桌宠'], [
    '产品形态：桌面悬浮宠物，点击穿透不遮挡操作，集提醒、命令、AI 对话于一体\n'
    'SOP 文档闭环：需求规范 → 技术规范 → 开发日志 → 验收清单，四大文档贯穿全流程\n'
    '工程化 SOP：一键打包脚本 + asar 纯净性校验 + 单测自动化，每次打包产出零个人数据产物\n'
    '安全 SOP：API Key 双重隔离（.env 不进仓库、不进安装包），渲染层零接触真实密钥',
])

# ============ 7. 核心业务模块（Slide 16 白底4项） ============
s = slides[15]
find_title(s).text_frame.text = '核心业务模块'
fill_body(s,
    ['提醒管理', '命令管理', '大模型对话', '交互系统'],
    [
        '定时调度引擎 10 秒轮询，到点触发气泡通知 + 音效，支持每日/一次性提醒',
        '常用命令看板支持搜索、置顶、一键复制，分类管理高频操作',
        'SSE 流式输出逐字渲染，思考过程分离展示，历史持久化可续聊',
        '单击互动动画、双击对话面板、右键配置中心，点击穿透 + 悬停恢复',
    ])

# ============ 8. 数据流与交互流（Slide 14 白底2项） ============
s = slides[13]
find_title(s).text_frame.text = '两条主线：数据流与交互流'
fill_body(s, ['数据流', '交互流'], [
    '渲染层操作 → IPC 通道 → 主进程 electron-store 落盘 → data:changed 广播 → 双窗口实时同步；\n对话数据经 chat.js 引擎流式回传',
    '窗口常驻穿透（forward 转发）→ 悬停命中可交互区动态恢复 → 单击/双击/右键判定（300ms 判定窗口）→\n拖动走 rAF 节流 + 绝对坐标锚点，跨屏不漂移',
])

# ============ 9. 技术架构（Slide 12 白底空白，自加分层文本框） ============
s = slides[11]
find_title(s).text_frame.text = '技术架构与关键技术'
arch = [
    ('渲染层（renderer）', '动画状态机 · 流式渲染缓冲 · 交互判定 · 拖动 rAF 节流'),
    ('preload 桥', 'contextBridge 白名单暴露 20+ 接口，contextIsolation 安全边界'),
    ('主进程（main）', '双窗口管理 · 10 秒提醒调度 · 置顶保活 · IPC 路由 · electron-store'),
    ('大模型后端（chat.js）', 'SSE 解析 · 思考分离 · 超窗摘要压缩 · 错误归一，纯 Node 零依赖'),
    ('DeepSeek API', 'deepseek-chat / deepseek-reasoner，Node 内置 fetch 直连'),
]
top = 1.75
for name, desc in arch:
    row = s.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(0.62))
    tf = row.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = name + '  '
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.name = '微软雅黑'
    r1.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(13)
    r2.font.name = '微软雅黑'
    r2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    top += 0.72

# ============ 10. 实施里程碑（Slide 17 白底4项） ============
s = slides[16]
find_title(s).text_frame.text = '实施里程碑'
fill_body(s,
    ['M0-M6 一期交付', 'v1.x 迭代', 'v2.x 迭代', 'v2.4 体验优化'],
    [
        '需求 → 原型 → UI → 交互 → 存储 → 通知 → 打包，七个里程碑全部通过验收',
        '配置中心、独立窗口、居中大窗、性能优化、NSIS 安装包发布',
        '交互面板改造、一键打包脚本、大模型流式会话接入',
        '置顶保活、拖动节流、动画减负——从「能用」到「好用」',
    ])

# ============ 11. 技术亮点（Slide 24 背景色4项） ============
s = slides[23]
find_title(s).text_frame.text = '技术亮点'
fill_body(s,
    ['点击穿透 + 动态恢复', '流式对话体验', '安全边界设计', '工程化 SOP 闭环'],
    [
        '穿透态不遮挡桌面操作，悬停命中宠物/面板自动恢复交互，体验无缝',
        '正文与思考分离渲染，超窗自动摘要压缩，错误归一友好提示',
        'Key 三重隔离（.env/存储/渲染层），打包白名单校验拒绝个人数据入包',
        '一键打包脚本 + 纯净校验 + 单测 + 文档规范，全流程可复现',
    ])

# ============ 12. 后续规划（Slide 21 背景色1项） ============
s = slides[20]
find_title(s).text_frame.text = '后续规划'
fill_body(s, ['从陪伴工具到智能助手'], [
    'Agent 能力：二期复刻 harness 的 agent 调度，让桌宠能操作桌面\n'
    '视觉升级：Live2D / 像素 GIF 素材替代 emoji，个性化皮肤体系\n'
    '轻量化：评估 Tauri 方案，安装包从 73MB 压缩到 10MB 级\n'
    '跨平台：macOS / Linux 适配，覆盖更多教学场景',
])

# ============ 13. 删除未使用的模板页（从后往前） ============
# 保留（1-based）：1封面 2扉页 5目录 8-17正文 21规划 24亮点 29谢谢
to_delete = [35, 34, 33, 32, 31, 30, 28, 27, 26, 25, 23, 22, 20, 19, 18, 7, 6, 4, 3]
to_delete_0 = [x - 1 for x in to_delete]
xml_slides = prs.slides._sldIdLst
ids = list(xml_slides)
for i in sorted(to_delete_0, reverse=True):
    sld = ids[i]
    rId = sld.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)

prs.save(OUT)
print(f'生成完成: {OUT}')
print(f'最终页数: {len(prs.slides._sldIdLst)}')
